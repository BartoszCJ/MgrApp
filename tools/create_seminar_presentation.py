from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from PIL import Image


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "output" / "pptx" / "prezentacja-seminarium-dyplomowe.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)


class C:
    BG = RGBColor(11, 16, 32)
    PANEL = RGBColor(23, 32, 51)
    PANEL2 = RGBColor(31, 41, 55)
    TEXT = RGBColor(248, 250, 252)
    MUTED = RGBColor(203, 213, 225)
    FAINT = RGBColor(148, 163, 184)
    LINE = RGBColor(51, 65, 85)
    BLUE = RGBColor(59, 130, 246)
    CYAN = RGBColor(6, 182, 212)
    GREEN = RGBColor(34, 197, 94)
    AMBER = RGBColor(245, 158, 11)
    RED = RGBColor(239, 68, 68)
    PURPLE = RGBColor(139, 92, 246)
    WHITE = RGBColor(255, 255, 255)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = C.BG
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C.BLUE
    bar.line.color.rgb = C.BLUE


def add_text(slide, text, x, y, w, h, size=14, color=C.TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None, n=None):
    add_text(slide, title, 0.55, 0.32, 8.8, 0.45, 26, C.TEXT, True)
    if subtitle:
        add_text(slide, subtitle, 0.56, 0.82, 9.4, 0.24, 9.5, C.FAINT)
    add_text(slide, "Bartosz Kloc | seminarium dyplomowe", 9.45, 0.42, 3.2, 0.2, 8.5, C.FAINT, False, PP_ALIGN.RIGHT)
    if n is not None:
        add_text(slide, f"{n:02d}", 12.35, 6.95, 0.45, 0.18, 8, C.FAINT, False, PP_ALIGN.RIGHT)


def add_bullets(slide, items, x, y, w, h, size=14):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = C.TEXT
        p.space_after = Pt(8)
    return box


def card(slide, x, y, w, h, heading, body, accent=C.BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C.PANEL
    shape.line.color.rgb = C.LINE
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.color.rgb = accent
    add_text(slide, heading, x + 0.22, y + 0.16, w - 0.36, 0.28, 14.5, C.TEXT, True)
    add_text(slide, body, x + 0.22, y + 0.55, w - 0.36, h - 0.65, 11.2, C.MUTED)


def pill(slide, x, y, text, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.95), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    add_text(slide, text, x, y + 0.11, 1.95, 0.18, 10.5, C.WHITE, True, PP_ALIGN.CENTER)


def fit_image(slide, image_path, x, y, w, h, caption=None):
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(5, 8, 22)
    frame.line.color.rgb = C.LINE
    path = str(image_path)
    with Image.open(path) as im:
        iw, ih = im.size
    box_ratio = (w - 0.16) / (h - 0.16)
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        pic_w = w - 0.16
        pic_h = pic_w / img_ratio
    else:
        pic_h = h - 0.16
        pic_w = pic_h * img_ratio
    px = x + 0.08 + ((w - 0.16) - pic_w) / 2
    py = y + 0.08 + ((h - 0.16) - pic_h) / 2
    slide.shapes.add_picture(path, Inches(px), Inches(py), Inches(pic_w), Inches(pic_h))
    if caption:
        add_text(slide, caption, x, y + h + 0.08, w, 0.18, 8.5, C.FAINT, False, PP_ALIGN.CENTER)


def metric(slide, x, y, value, label, color):
    add_text(slide, value, x, y, 1.4, 0.5, 27, color, True, PP_ALIGN.CENTER)
    add_text(slide, label, x - 0.15, y + 0.58, 1.7, 0.4, 8.5, C.MUTED, False, PP_ALIGN.CENTER)


def arrow(slide, x, y, w=0.55):
    line = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    line.line.color.rgb = C.FAINT
    line.line.width = Pt(2)


images = {
    "ui": ROOT / "tmp" / "CoJest-media" / "image1.png",
    "graph": ROOT / "tmp" / "CoJest-media" / "image2.png",
    "results": ROOT / "tmp" / "CoJest-media" / "image5.png",
}


slides = []
for _ in range(14):
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    slides.append(s)

s = slides[0]
add_text(s, "Analiza aktywności\ntransakcyjnej w blockchainie", 0.7, 1.0, 7.2, 1.45, 34, C.TEXT, True)
add_text(s, "po incydentach cyberbezpieczeństwa", 0.72, 2.55, 7.0, 0.36, 20, C.CYAN, True)
add_text(s, "Seminarium dyplomowe | Bartosz Kloc | promotor: dr Ruslan Shevchuk", 0.73, 3.18, 7.2, 0.28, 13, C.MUTED)
pill(s, 0.75, 4.15, "Ronin", C.BLUE)
pill(s, 2.9, 4.15, "Euler", C.PURPLE)
pill(s, 5.05, 4.15, "Nomad", C.GREEN)
circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.7), Inches(1.1), Inches(3.8), Inches(3.8))
circle.fill.solid()
circle.fill.fore_color.rgb = C.PANEL
circle.line.color.rgb = C.BLUE
circle.line.width = Pt(3)
add_text(s, "prototyp\nanaliza grafowa\nwyniki", 9.4, 2.15, 2.55, 0.9, 18, C.TEXT, True, PP_ALIGN.CENTER)
add_text(s, "01", 12.35, 6.95, 0.45, 0.18, 8, C.FAINT, False, PP_ALIGN.RIGHT)

s = slides[1]
add_title(s, "Plan prezentacji", "Od problemu do aktualnych wyników prototypu", 2)
for i, (n, h, b) in enumerate([
    ("1", "Motywacja i problem", "Dlaczego śledzenie przepływu środków po incydencie jest trudne."),
    ("2", "Cel i zakres", "Co obejmuje praca i jakie pytania ma rozwiązać."),
    ("3", "Rozwiązanie", "Architektura, metoda śledzenia i implementacja."),
    ("4", "Eksperymenty i wnioski", "Wyniki dla Ronin, Euler i Nomad oraz dalsze prace."),
]):
    y = 1.45 + i * 1.12
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), Inches(y), Inches(0.56), Inches(0.56))
    o.fill.solid(); o.fill.fore_color.rgb = C.BLUE; o.line.color.rgb = C.BLUE
    add_text(s, n, 0.85, y + 0.17, 0.56, 0.14, 12, C.WHITE, True, PP_ALIGN.CENTER)
    add_text(s, h, 1.65, y - 0.02, 4.0, 0.26, 17, C.TEXT, True)
    add_text(s, b, 1.65, y + 0.35, 9.4, 0.28, 12.5, C.MUTED)

s = slides[2]
add_title(s, "Wprowadzenie i motywacja", "Przykład: po ataku środki szybko trafiają na wiele adresów", 3)
add_text(s, "Ronin Bridge", 0.75, 1.35, 3, 0.32, 22, C.TEXT, True)
add_text(s, "$625M", 0.75, 1.85, 3.4, 0.65, 46, C.AMBER, True)
add_text(s, "jeden z największych incydentów blockchainowych", 0.78, 2.62, 3.7, 0.35, 12.5, C.MUTED)
card(s, 5.05, 1.35, 2.3, 1.3, "Adres źródłowy", "punkt startowy analizy po incydencie", C.RED)
card(s, 8.0, 0.85, 2.3, 1.05, "Pośrednicy", "kolejne adresy i transfery", C.BLUE)
card(s, 8.0, 2.25, 2.3, 1.05, "Punkty końcowe", "giełdy, mosty, miksery", C.GREEN)
card(s, 10.75, 1.55, 1.85, 1.25, "Wynik", "graf i metryki", C.CYAN)
add_bullets(s, ["ręczna analiza w eksploratorze jest czasochłonna", "środki mogą przechodzić przez wiele adresów", "potrzebne jest narzędzie porządkujące analizę"], 0.8, 4.15, 6.2, 1.4, 15)

s = slides[3]
add_title(s, "Problem badawczy i cel pracy", "Automatyzacja części analizy po znanym incydencie", 4)
add_text(s, "Czy można zautomatyzować analizę przepływu środków po incydencie?", 0.8, 1.25, 11.5, 0.75, 23, C.TEXT, True)
card(s, 0.8, 2.45, 3.7, 1.45, "Odtworzenie przepływu", "budowa grafu transakcyjnego od adresu startowego", C.BLUE)
card(s, 4.8, 2.45, 3.7, 1.45, "Wykrycie wzorców", "reguły dla CEX, Tornado Cash, bridge i peel chain", C.AMBER)
card(s, 8.8, 2.45, 3.7, 1.45, "Ocena jakości", "porównanie wyniku z ground truth dla znanych przypadków", C.GREEN)
add_text(s, "Cel pracy: opracowanie prototypu wspierającego analizę aktywności transakcyjnej po incydentach w sieci Ethereum.", 1.0, 4.75, 11.2, 0.8, 17, C.MUTED, False, PP_ALIGN.CENTER)

s = slides[4]
add_title(s, "Zakres i zadania badawcze", "Zakres zaakceptowany do dalszej pracy: Ethereum + 3 case studies", 5)
card(s, 0.75, 1.35, 2.7, 1.2, "Zakres", "Ethereum / EVM\nanaliza po incydencie", C.BLUE)
card(s, 3.75, 1.35, 2.7, 1.2, "Przypadki", "Ronin Bridge\nEuler Finance\nNomad Bridge", C.PURPLE)
card(s, 6.75, 1.35, 2.7, 1.2, "Dane", "Etherscan\nArkham\nopisy incydentów", C.GREEN)
card(s, 9.75, 1.35, 2.7, 1.2, "Poza zakresem", "Bitcoin\ninne sieci\npełne smart calls", C.RED)
add_text(s, "Główne zadania badawcze", 0.85, 3.35, 4.3, 0.3, 18, C.TEXT, True)
add_bullets(s, ["opracowanie modelu śledzenia przepływu środków", "implementacja prototypu aplikacji", "dobór heurystyk i metryk oceny", "przeprowadzenie eksperymentów na znanych incydentach"], 0.95, 3.9, 8.4, 1.8, 14.5)

s = slides[5]
add_title(s, "Tło problemu", "Istniejące narzędzia pomagają, ale praca skupia się na prototypie badawczym", 6)
for i, (h, b, c) in enumerate([
    ("Eksploratory blockchain", "Pokazują transakcje, ale analiza wielu kroków jest ręczna.", C.BLUE),
    ("Narzędzia komercyjne", "Chainalysis / Elliptic są pomocne, ale zamknięte badawczo.", C.PURPLE),
    ("Metody ML/GNN", "Osobny kierunek badań, wymagający dużych zbiorów danych.", C.CYAN),
    ("Moje miejsce", "Analiza po incydencie: graf przepływu, heurystyki i wyniki.", C.GREEN),
]):
    card(s, 0.85 + (i % 2) * 5.95, 1.35 + (i // 2) * 2.15, 5.4, 1.55, h, b, c)

s = slides[6]
add_title(s, "Proponowane rozwiązanie", "Architektura prototypu", 7)
for i, (h, b, c) in enumerate([
    ("Frontend", "Next.js\ninterfejs użytkownika", C.BLUE),
    ("Backend", "FastAPI\nlogika analizy", C.PURPLE),
    ("Źródła danych", "Etherscan\nArkham", C.CYAN),
    ("Cache", "lokalny zapis\nodpowiedzi API", C.AMBER),
    ("Analiza", "graf + heurystyki\nmetryki", C.GREEN),
]):
    x = 0.65 + i * 2.2
    card(s, x, 2.1, 1.75, 1.45, h, b, c)
    if i < 4:
        arrow(s, x + 1.84, 2.82, 0.28)
add_text(s, "Efekt: użytkownik może uruchomić analizę, obejrzeć graf przepływu i porównać wyniki dla wybranych przypadków.", 1.0, 4.65, 11.2, 0.5, 17, C.MUTED, False, PP_ALIGN.CENTER)

s = slides[7]
add_title(s, "Metoda śledzenia przepływu", "BFS, czyli analiza kolejnych kroków od adresu startowego", 8)
fit_image(s, images["graph"], 6.0, 1.25, 6.5, 4.55, "Przykładowy widok grafu przepływu w aplikacji")
add_bullets(s, ["start od adresu powiązanego z incydentem", "kolejne transakcje tworzą graf przepływu", "hops oznacza głębokość śledzenia", "CEX, mikser lub bridge to istotne punkty końcowe"], 0.75, 1.55, 4.55, 3.0, 14)
add_text(s, "Główne ustawienie wyników: hops=3.", 0.78, 5.25, 4.4, 0.35, 13.5, C.AMBER, True)

s = slides[8]
add_title(s, "Implementacja prototypu", "Działająca aplikacja webowa z widokiem grafu i wyników", 9)
fit_image(s, images["ui"], 0.75, 1.25, 6.0, 4.55, "Ekran główny aplikacji")
fit_image(s, images["results"], 7.05, 1.25, 5.55, 3.0, "Tabela eksperymentów dla 3 przypadków")
add_text(s, "Zaimplementowane elementy", 7.1, 4.55, 3.7, 0.25, 15, C.TEXT, True)
add_bullets(s, ["frontend i backend", "integracja z API", "cache odpowiedzi", "graf, heurystyki i metryki"], 7.15, 4.95, 4.6, 1.25, 12.5)

s = slides[9]
add_title(s, "Eksperymenty", "Trzy znane incydenty jako studia przypadków", 10)
for i, (h, b, c) in enumerate([
    ("Ronin Bridge", "$625M, 2022\nnajważniejszy przypadek dla CEX Coverage", C.BLUE),
    ("Euler Finance", "$197M, 2023\nczęściowe odtworzenie znanych adresów", C.PURPLE),
    ("Nomad Bridge", "$190M, 2022\nnajwyższy Address Recall w eksperymencie", C.GREEN),
]):
    card(s, 0.85 + i * 4.1, 1.35, 3.6, 1.55, h, b, c)
add_text(s, "Metryki oceny", 0.9, 3.65, 3.0, 0.3, 18, C.TEXT, True)
for i, (h, b, c) in enumerate([
    ("Address Recall", "ile znanych adresów znaleziono", C.BLUE),
    ("Heuristic Precision", "ile alertów było uzasadnionych", C.AMBER),
    ("Heuristic Recall", "ile oczekiwanych kategorii wykryto", C.GREEN),
    ("CEX Coverage", "czy dotarto do znanych adresów giełdowych", C.CYAN),
]):
    card(s, 0.9 + (i % 2) * 5.75, 4.15 + (i // 2) * 0.85, 5.2, 0.65, h, b, c)

s = slides[10]
add_title(s, "Wyniki i interpretacja", "Finalny eksperyment: hops=3", 11)
rows = [
    ["Przypadek", "Addr. Recall", "Heur. Precision", "Heur. Recall", "CEX Coverage"],
    ["Ronin", "70%", "67%", "67%", "100%"],
    ["Euler", "50%", "33%", "50%", "N/A"],
    ["Nomad", "88%", "67%", "67%", "N/A"],
    ["Średnia", "69%", "56%", "61%", "100%"],
]
table = s.shapes.add_table(len(rows), len(rows[0]), Inches(0.75), Inches(1.35), Inches(7.25), Inches(2.9)).table
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = C.PANEL if r else C.PANEL2
        for p in cell.text_frame.paragraphs:
            p.font.name = "Aptos"
            p.font.size = Pt(10.5)
            p.font.color.rgb = C.TEXT
            if r in (0, 4):
                p.font.bold = True
metric(s, 8.55, 1.45, "88%", "najwyższy Address Recall\nNomad", C.GREEN)
metric(s, 10.2, 1.45, "100%", "CEX Coverage\nRonin", C.GREEN)
metric(s, 11.85, 1.45, "3", "studia\nprzypadków", C.BLUE)
add_bullets(s, ["Ronin: dotarto do adresów giełdowych", "Nomad: najwyższy zasięg adresów", "hops=3 daje bogatszy graf, ale więcej zapytań"], 8.45, 3.35, 3.65, 1.8, 12)

s = slides[11]
add_title(s, "Ograniczenia i dyskusja", "Uczciwy zakres prototypu", 12)
for i, (h, b, c) in enumerate([
    ("Zakres sieci", "obecnie Ethereum / EVM, bez Bitcoina i innych modeli transakcyjnych", C.BLUE),
    ("Jakość danych", "wyniki zależą od publicznych API, etykiet i udokumentowanego ground truth", C.AMBER),
    ("Głębokość analizy", "większy hops zwiększa skuteczność, ale szybko powiększa graf", C.GREEN),
    ("Smart contracts", "pełna analiza wewnętrznych wywołań kontraktów jest poza aktualnym zakresem", C.RED),
]):
    card(s, 0.85 + (i % 2) * 5.9, 1.35 + (i // 2) * 1.75, 5.35, 1.15, h, b, c)
add_text(s, "Te ograniczenia nie blokują pracy — wyznaczają granice prototypu i naturalne kierunki rozwoju.", 1.0, 5.35, 11.2, 0.45, 16, C.MUTED, False, PP_ALIGN.CENTER)

s = slides[12]
add_title(s, "Wnioski i dalsze prace", "Aktualny stan: prototyp działa, teraz praca pisemna", 13)
card(s, 0.8, 1.3, 3.65, 1.35, "Wniosek 1", "prototyp działa end-to-end", C.GREEN)
card(s, 4.85, 1.3, 3.65, 1.35, "Wniosek 2", "wyniki można porównać z ground truth", C.BLUE)
card(s, 8.9, 1.3, 3.65, 1.35, "Wniosek 3", "hops=3 poprawia zasięg, ale zwiększa koszt", C.AMBER)
add_text(s, "Dalsze kroki", 0.9, 3.45, 2.7, 0.28, 18, C.TEXT, True)
add_bullets(s, ["opis części teoretycznej: blockchain, Ethereum, OSINT i forensics", "opis metodologii: graf, heurystyki, metryki i ground truth", "interpretacja wyników eksperymentów", "w przyszłości: inne sieci, większa baza incydentów i dokładniejsza analiza smart kontraktów"], 0.95, 3.9, 9.8, 1.7, 14.5)

s = slides[13]
add_text(s, "Dziękuję za uwagę", 0.8, 2.2, 11.8, 0.65, 42, C.TEXT, True, PP_ALIGN.CENTER)
add_text(s, "Pytania?", 0.8, 3.15, 11.8, 0.45, 24, C.CYAN, True, PP_ALIGN.CENTER)
add_text(s, "Automatyczne wykrywanie i analiza aktywności transakcyjnej w blockchainie po incydentach cyberbezpieczeństwa", 1.8, 4.45, 9.75, 0.45, 13, C.MUTED, False, PP_ALIGN.CENTER)
add_text(s, "14", 12.35, 6.95, 0.45, 0.18, 8, C.FAINT, False, PP_ALIGN.RIGHT)

prs.save(OUT)
print(OUT)
